import dash
from dash import dcc, html, dash_table
import dash_bootstrap_components as dbc
import plotly.graph_objects as go
import plotly.express as px

import threading
import pandas as pd
import datetime
import ipaddress
import os
import sys
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# --- Optional / fragile dependencies are imported defensively -------------
# scapy needs raw-socket privileges (root/admin) to sniff or send packets.
# python-nmap needs the *nmap* binary installed on the OS, not just the pip
# package. Neither of these should be able to stop the Dash server from
# starting, so we guard them and just disable the related features if
# something is missing.

try:
    from scapy.all import sniff, IP, sr1, TCP
    SCAPY_AVAILABLE = True
except Exception as e:
    print(f"[WARN] scapy not usable ({e}). Packet sniffing/scanning disabled.")
    SCAPY_AVAILABLE = False

try:
    import nmap  # python-nmap library
    nm = nmap.PortScanner()  # this line throws if the `nmap` binary isn't on PATH
    NMAP_AVAILABLE = True
except Exception as e:
    print(f"[WARN] nmap not usable ({e}). Install the 'nmap' program "
          f"(not just pip install python-nmap) to enable detailed scans.")
    nm = None
    NMAP_AVAILABLE = False

# Initialize the Dash app
app = dash.Dash(__name__, external_stylesheets=[dbc.themes.CYBORG])
server = app.server  # For deployment

# Data storage
packet_data = []
external_scan_results = []
scan_progress = ""

# IP Classification function
def is_internal(ip):
    internal_ranges = [
        '10.0.0.0/8',
        '172.16.0.0/12',
        '192.168.0.0/16',
        '127.0.0.0/8'
    ]
    try:
        ip_obj = ipaddress.IPv4Address(ip)
        for network in internal_ranges:
            if ip_obj in ipaddress.IPv4Network(network):
                return True
        return False
    except Exception:
        return False

# Packet capture callback
def packet_callback(packet):
    global packet_data
    if packet.haslayer(IP):
        src = packet[IP].src
        dst = packet[IP].dst
        new_entry = {
            "Timestamp": datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "Source IP": src,
            "Destination IP": dst,
            "Protocol": packet[IP].proto,
            "Length": len(packet),
            "Source External": not is_internal(src),
            "Destination External": not is_internal(dst)
        }
        packet_data.append(new_entry)
        if len(packet_data) > 1000:
            packet_data.pop(0)

# Start packet sniffing in a separate thread (guarded so a permissions
# error doesn't just vanish into an unhandled thread exception)
def start_sniffing():
    try:
        sniff(prn=packet_callback, store=False)
    except PermissionError:
        print("[ERROR] Packet sniffing needs elevated privileges. "
              "Run this script with sudo (Linux/Mac) or as Administrator "
              "with Npcap installed (Windows). Sniffing disabled for now.")
    except Exception as e:
        print(f"[ERROR] Sniffing thread failed: {e}")

if SCAPY_AVAILABLE:
    threading.Thread(target=start_sniffing, daemon=True).start()

# Lightweight port scanner using scapy
def scapy_port_scan(ip, ports=[80]):
    global external_scan_results, scan_progress
    if not SCAPY_AVAILABLE:
        scan_progress = "Scapy not available on this system"
        return []
    results = []
    for port in ports:
        scan_progress = f"Scanning {ip}:{port}"
        packet = IP(dst=ip)/TCP(dport=port, flags="S")
        response = sr1(packet, timeout=1, verbose=0)
        if response and response.haslayer(TCP):
            if response[TCP].flags == 0x12:  # SYN-ACK
                results.append({"IP": ip, "Port": port, "Status": "Open"})
                external_scan_results.append({"IP": ip, "Port": port})
            elif response[TCP].flags == 0x14:  # RST-ACK
                results.append({"IP": ip, "Port": port, "Status": "Closed"})
        else:
            results.append({"IP": ip, "Port": port, "Status": "Filtered"})
    scan_progress = "Scan completed"
    return results

# Nmap-based scanning
def nmap_scan(targets="127.0.0.1", ports="80"):
    global external_scan_results, scan_progress
    if not NMAP_AVAILABLE:
        scan_progress = "nmap binary not installed on this system"
        return False
    try:
        scan_progress = "Starting scan..."
        nm.scan(hosts=targets, ports=ports, arguments='-T4')
        scan_progress = "Scan completed"

        for host in nm.all_hosts():
            for proto in nm[host].all_protocols():
                port_list = nm[host][proto].keys()
                for port in port_list:
                    if nm[host][proto][port]['state'] == 'open':
                        external_scan_results.append({
                            "IP": host,
                            "Port": port,
                            "Service": nm[host][proto][port]['name']
                        })
        return True
    except Exception as e:
        scan_progress = f"Scan error: {str(e)}"
        return False

# Generate PDF report
def generate_pdf_report(df):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Intrusion Detection Report", ln=True, align="C")
    pdf.ln(10)
    for index, row in df.iterrows():
        pdf.cell(200, 10, txt=f"Timestamp: {row['Timestamp']}", ln=True)
        pdf.cell(200, 10, txt=f"Source IP: {row['Source IP']}", ln=True)
        pdf.cell(200, 10, txt=f"Destination IP: {row['Destination IP']}", ln=True)
        pdf.cell(200, 10, txt=f"Protocol: {row['Protocol']}", ln=True)
        pdf.cell(200, 10, txt=f"Packet Length: {row['Length']}", ln=True)
        pdf.cell(200, 10, txt=f"Source External: {row['Source External']}", ln=True)
        pdf.cell(200, 10, txt=f"Destination External: {row['Destination External']}", ln=True)
        pdf.ln(10)
    pdf_filename = "intrusion_report.pdf"
    pdf.output(pdf_filename)
    return pdf_filename

# Generate PowerPoint report
def generate_ppt_report(df):
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Intrusion Detection Report"
    subtitle.text = "Generated on: " + datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    # Summary Slide
    summary_slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(summary_slide_layout)
    title = slide.shapes.title
    title.text = "Summary"
    left = top = Inches(1)
    width = height = Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = (
        f"Total Packets Captured: {len(df)}\n"
        f"Unique Source IPs: {df['Source IP'].nunique()}\n"
        f"Unique Destination IPs: {df['Destination IP'].nunique()}\n"
        f"Protocols Detected: {', '.join(map(str, df['Protocol'].unique()))}"
    )

    # Protocol Distribution Chart
    protocol_counts = df['Protocol'].value_counts().to_dict()
    fig_protocol = go.Figure(go.Pie(
        labels=list(protocol_counts.keys()),
        values=list(protocol_counts.values()),
        hole=0.3,
        marker=dict(colors=px.colors.sequential.RdBu)
    ))
    fig_protocol.update_layout(
        title="Protocol Distribution",
        plot_bgcolor='#222',
        paper_bgcolor='#222',
        font=dict(color='white')
    )
    protocol_image = BytesIO(fig_protocol.to_image(format="png"))
    slide = prs.slides.add_slide(summary_slide_layout)
    slide.shapes.title.text = "Protocol Distribution"
    slide.shapes.add_picture(protocol_image, Inches(1), Inches(2), width=Inches(6))

    # Save the PPT
    ppt_filename = "intrusion_report.pptx"
    prs.save(ppt_filename)
    return ppt_filename

# Dashboard layout
app.layout = dbc.Container([
    dbc.NavbarSimple(
        brand="Intrusion Detection Dashboard",
        color="primary",
        dark=True,
        className="mb-3"
    ),
    dbc.Row([
        dbc.Col([
            dbc.Card([
                dbc.CardHeader("Live Packet Monitoring"),
                dbc.CardBody(dcc.Graph(id='live-packet-chart'))
            ]),
        ], width=6),
        dbc.Col([
            dbc.Card([
                dbc.CardHeader("Protocol Breakdown"),
                dbc.CardBody(dcc.Graph(id='protocol-chart'))
            ]),
        ], width=6)
    ]),
    dbc.Row([
        dbc.Col([
            dbc.Card([
                dbc.CardHeader("External Traffic Analysis"),
                dbc.CardBody(dcc.Graph(id='external-traffic-chart'))
            ]),
        ], width=12)
    ]),
    dbc.Row([
        dbc.Col([
            dbc.Button("Generate Report", id="report-btn", color="success", className="mb-3"),
            dcc.Download(id="download-csv"),
            dcc.Download(id="download-pdf"),
            dcc.Download(id="download-ppt"),
            dbc.Button("Scan External IPs", id="scan-btn", color="warning", className="ml-3 mb-3"),
            dbc.Input(id="scan-target", placeholder="Enter IP or range (e.g., 192.168.1.1 or 192.168.1.0/24)",
                     type="text", className="mb-2"),
            dbc.Input(id="scan-ports", placeholder="Enter ports (e.g., 80,443 or 1-100)",
                     type="text", value="80", className="mb-2"),
            dbc.RadioItems(
                id="scan-method",
                options=[
                    {"label": "Quick Scan (Scapy)", "value": "scapy"},
                    {"label": "Detailed Scan (Nmap)", "value": "nmap"}
                ],
                value="scapy",
                inline=True,
                className="mb-2"
            ),
            html.Div(id="scan-progress", className="text-info mb-2"),
            html.Div(id="scan-results", className="mt-3")
        ], width=12)
    ]),
    dbc.Row([
        dbc.Col([
            dbc.Card([
                dbc.CardHeader("Packet Details Table"),
                dbc.CardBody(dash_table.DataTable(
                    id='packet-table',
                    style_table={'overflowX': 'auto'},
                    style_cell={
                        'textAlign': 'center',
                        'padding': '5px',
                        'backgroundColor': '#222',
                        'color': 'white',
                        'maxWidth': '150px'
                    },
                    style_header={
                        'backgroundColor': '#444',
                        'color': 'white',
                        'fontWeight': 'bold'
                    }
                ))
            ])
        ], width=12)
    ]),
    dcc.Interval(id='interval-update', interval=2000, n_intervals=0)
])

# Dashboard update callback
@app.callback(
    [
        dash.Output('live-packet-chart', 'figure'),
        dash.Output('protocol-chart', 'figure'),
        dash.Output('external-traffic-chart', 'figure'),
        dash.Output('packet-table', 'data'),
        dash.Output('packet-table', 'columns')
    ],
    [dash.Input('interval-update', 'n_intervals')]
)
def update_dashboard(n):
    global packet_data
    if not packet_data:
        return go.Figure(), go.Figure(), go.Figure(), [], []
    df = pd.DataFrame(packet_data)

    # Line Chart for Packet Monitoring
    fig_packet = go.Figure()
    fig_packet.add_trace(go.Scatter(
        x=df["Timestamp"],
        y=df["Length"],
        mode='lines+markers',
        name='Packet Length',
        marker=dict(color='#EF553B')
    ))
    fig_packet.update_layout(
        title="Live Packet Monitoring",
        xaxis_title="Time",
        yaxis_title="Packet Size",
        template="plotly_dark",
        plot_bgcolor='#222',
        paper_bgcolor='#222',
        font=dict(color='white')
    )

    # Protocol Pie Chart
    protocol_counts = df['Protocol'].value_counts().to_dict()
    fig_protocol = go.Figure(go.Pie(
        labels=list(protocol_counts.keys()),
        values=list(protocol_counts.values()),
        hole=0.3,
        marker=dict(colors=px.colors.sequential.RdBu)
    ))
    fig_protocol.update_layout(
        title="Protocol Distribution",
        plot_bgcolor='#222',
        paper_bgcolor='#222',
        font=dict(color='white')
    )

    # External Traffic Bar Chart
    external_traffic = df.groupby(['Source External', 'Destination External']).size().reset_index(name='Count')
    external_traffic['Category'] = external_traffic.apply(
        lambda row: f"{'External' if row['Source External'] else 'Internal'} → {'External' if row['Destination External'] else 'Internal'}",
        axis=1
    )
    fig_external = px.bar(
        external_traffic,
        x='Category',
        y='Count',
        color='Count',
        color_continuous_scale='Teal',
        title='Internal/External Traffic Patterns'
    )
    fig_external.update_layout(
        plot_bgcolor='#222',
        paper_bgcolor='#222',
        font=dict(color='white'),
        xaxis=dict(tickangle=45)
    )

    # Prepare table data and columns
    columns = [
        {"name": "Timestamp", "id": "Timestamp"},
        {"name": "Source IP", "id": "Source IP"},
        {"name": "Destination IP", "id": "Destination IP"},
        {"name": "Protocol", "id": "Protocol"},
        {"name": "Length", "id": "Length"},
        {"name": "Src External", "id": "Source External"},
        {"name": "Dst External", "id": "Destination External"}
    ]
    table_data = df.to_dict('records')
    return fig_packet, fig_protocol, fig_external, table_data, columns

# Report generation callback
@app.callback(
    [
        dash.Output("download-csv", "data"),
        dash.Output("download-pdf", "data"),
        dash.Output("download-ppt", "data")
    ],
    [dash.Input("report-btn", "n_clicks")],
    prevent_initial_call=True
)
def generate_report(n_clicks):
    df = pd.DataFrame(packet_data)
    if df.empty:
        empty = dcc.send_bytes(b"No packet data captured yet.", "empty_report.txt")
        return empty, empty, empty

    csv_filename = pdf_filename = ppt_filename = None
    try:
        csv_filename = "intrusion_report.csv"
        df.to_csv(csv_filename, index=False)
        csv_data = dcc.send_file(csv_filename)
    except Exception as e:
        print(f"[ERROR] CSV generation failed: {e}")
        csv_data = dcc.send_bytes(str(e).encode(), "csv_error.txt")

    try:
        pdf_filename = generate_pdf_report(df)
        pdf_data = dcc.send_file(pdf_filename)
    except Exception as e:
        print(f"[ERROR] PDF generation failed: {e}")
        pdf_data = dcc.send_bytes(str(e).encode(), "pdf_error.txt")

    try:
        ppt_filename = generate_ppt_report(df)
        ppt_data = dcc.send_file(ppt_filename)
    except Exception as e:
        print(f"[ERROR] PPT generation failed: {e} "
              f"(if this mentions kaleido, run: pip install -U kaleido)")
        ppt_data = dcc.send_bytes(str(e).encode(), "ppt_error.txt")

    for f in (csv_filename, pdf_filename, ppt_filename):
        if f and os.path.exists(f):
            os.remove(f)

    return csv_data, pdf_data, ppt_data

# External IP scan callback
@app.callback(
    [
        dash.Output("scan-results", "children"),
        dash.Output("scan-progress", "children")
    ],
    [dash.Input("scan-btn", "n_clicks")],
    [
        dash.State("scan-target", "value"),
        dash.State("scan-ports", "value"),
        dash.State("scan-method", "value")
    ],
    prevent_initial_call=True
)
def scan_external_ips(n_clicks, target, ports, method):
    global external_scan_results, scan_progress
    external_scan_results.clear()

    if not target:
        return "Please enter a target IP or range", "Error: No target specified"

    ports = ports if ports else "80"

    try:
        if method == "scapy":
            scan_progress = "Starting quick scan with Scapy..."
            target_ip = target.split('/')[0] if '/' in target else target
            port_list = [int(p) for p in ports.split(',')] if ',' in ports else [int(ports)]
            results = scapy_port_scan(target_ip, port_list)

            if not results:
                return "No open ports found (or scapy unavailable — see terminal)", "Scan completed - no open ports"

            scan_df = pd.DataFrame(results)
            return [
                dash_table.DataTable(
                    data=scan_df.to_dict('records'),
                    columns=[{"name": i, "id": i} for i in scan_df.columns],
                    style_table={'overflowX': 'auto'},
                    style_cell={
                        'textAlign': 'center',
                        'padding': '5px',
                        'backgroundColor': '#222',
                        'color': 'white',
                        'maxWidth': '150px'
                    },
                    style_header={
                        'backgroundColor': '#444',
                        'color': 'white',
                        'fontWeight': 'bold'
                    }
                )
            ], "Scan completed"

        elif method == "nmap":
            scan_progress = "Starting detailed scan with Nmap..."
            nmap_scan(target, ports)

            if not external_scan_results:
                return "No open ports found (or nmap unavailable — see terminal)", "Scan completed - no open ports"

            scan_df = pd.DataFrame(external_scan_results)
            return [
                dash_table.DataTable(
                    data=scan_df.to_dict('records'),
                    columns=[{"name": i, "id": i} for i in scan_df.columns],
                    style_table={'overflowX': 'auto'},
                    style_cell={
                        'textAlign': 'center',
                        'padding': '5px',
                        'backgroundColor': '#222',
                        'color': 'white',
                        'maxWidth': '150px'
                    },
                    style_header={
                        'backgroundColor': '#444',
                        'color': 'white',
                        'fontWeight': 'bold'
                    }
                )
            ], "Scan completed"

    except Exception as e:
        return f"Scan error: {str(e)}", f"Error: {str(e)}"

if __name__ == '__main__':
    port = 8050
    print("=" * 60)
    print(f"  Dashboard starting -> open this in Chrome:")
    print(f"  http://127.0.0.1:{port}")
    print("=" * 60)
    # host='0.0.0.0' so it's reachable if you're in a container / WSL /
    # remote-SSH / Codespaces setup, where 127.0.0.1 inside VS Code's
    # terminal is not the same machine as the Chrome window.
    # If you're running fully locally (VS Code + Chrome on the same
    # physical machine), 127.0.0.1 also works fine.
    app.run(debug=True, host='0.0.0.0', port=port)